"""Office and PPT workflow helpers for the single-agent runtime."""

from __future__ import annotations

import re
import shutil
import time
from pathlib import Path

from whaleclaw.providers.base import Message
from whaleclaw.sessions.manager import Session

OFFICE_PATH_RE = re.compile(r"(/[^\n\"']+\.(?:pptx|docx|xlsx))", re.IGNORECASE)
ABS_FILE_PATH_RE = re.compile(r"(/[^\s\"')]+?\.[A-Za-z0-9]{1,8})(?=[\s\"')]|$)")
NON_DELIVERY_EXTS = {".py", ".sh", ".bash", ".zsh", ".log", ".tmp"}
VERSION_SUFFIX_RE = re.compile(r"_V\d+$", re.IGNORECASE)


def is_office_path_probe_command(command: str) -> bool:
    low = command.strip().lower()
    if not low:
        return False
    probe_tools = (
        "find ",
        "mdfind ",
        "locate ",
        "fd ",
        "ls ",
        "stat ",
        "test -f ",
        "test -e ",
    )
    office_terms = (
        "ppt",
        "pptx",
        "docx",
        "xlsx",
        "word",
        "excel",
        "幻灯片",
        "文档",
        "表格",
        ".pptx",
        ".docx",
        ".xlsx",
    )
    return any(token in low for token in probe_tools) and any(
        token in low for token in office_terms
    )


def looks_like_ppt_generation_script(text: str) -> bool:
    low = text.lower()
    if not low:
        return False
    hints = (
        "from pptx import presentation",
        "pptx import presentation",
        "presentation()",
        ".save(",
        ".pptx",
    )
    return any(hint in low for hint in hints)


def looks_like_ppt_generation_command(command: str) -> bool:
    low = command.lower()
    if not low:
        return False
    if "python" not in low and "python3" not in low:
        return False
    return ".pptx" in low or "pptx" in low or "presentation(" in low


def extract_office_paths(text: str) -> list[str]:
    if not text:
        return []
    paths: list[str] = []
    for match in OFFICE_PATH_RE.finditer(text):
        path = match.group(1).strip()
        if path:
            paths.append(path)
    return paths


def extract_round_delivery_section(text: str) -> str:
    if not text:
        return ""
    marker = re.search(r"(?m)^\s*[1-5][\)\.、:：]\s*本轮可直接交付结果", text)
    if marker is None:
        marker = re.search(r"(?m)^\s*[1-5][\)\.、:：]\s*本轮可直接使用", text)
    if marker is None:
        return text
    tail = text[marker.start() :]
    next_idx = len(tail)
    next_marker = re.search(r"(?m)^\s*[2-9][\)\.、:：](?!\s*本轮可直接)", tail)
    if next_marker is not None and next_marker.start() > 0:
        next_idx = next_marker.start()
    return tail[:next_idx].strip()


def extract_delivery_artifact_paths(
    text: str,
    *,
    include_scripts: bool = False,
) -> list[str]:
    section = extract_round_delivery_section(text)
    if not section:
        return []
    output: list[str] = []
    seen: set[str] = set()
    for match in ABS_FILE_PATH_RE.finditer(section):
        path = match.group(1).strip()
        if not path:
            continue
        suffix = Path(path).suffix.lower()
        if (not include_scripts) and suffix in NON_DELIVERY_EXTS:
            continue
        if path in seen:
            continue
        seen.add(path)
        output.append(path)
    return output


def with_round_version_suffix(path: str, round_no: int) -> str:
    raw = str(Path(path).expanduser())
    target = Path(raw)
    stem = target.stem
    base_stem = VERSION_SUFFIX_RE.sub("", stem)
    target_name = f"{base_stem}_V{round_no}{target.suffix}"
    return str(target.with_name(target_name))


def fix_version_suffix(paths: list[str], round_no: int) -> tuple[list[str], dict[str, str]]:
    fixed: list[str] = []
    rename_map: dict[str, str] = {}
    expected = f"_V{round_no}"
    for path in paths:
        file_path = Path(path)
        stem = file_path.stem
        match = VERSION_SUFFIX_RE.search(stem)
        if match is not None and match.group(0).upper() != expected.upper():
            correct = with_round_version_suffix(path, round_no)
            correct_path = Path(correct)
            try:
                correct_path.parent.mkdir(parents=True, exist_ok=True)
                if correct_path.exists():
                    correct_path.unlink()
                shutil.copy2(file_path, correct_path)
                rename_map[path] = str(correct_path)
                fixed.append(str(correct_path))
                continue
            except OSError:
                pass
        fixed.append(path)
    return fixed, rename_map


def extract_artifact_baseline(paths: list[str]) -> str:
    for path in paths:
        file_path = Path(path).expanduser()
        if not file_path.exists():
            continue
        if file_path.suffix.lower() == ".pptx":
            return baseline_from_pptx(file_path)
    return ""


def baseline_from_pptx(file_path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return ""
    try:
        presentation = Presentation(str(file_path))
    except Exception:
        return ""
    slide_count = len(presentation.slides)
    total_images = 0
    total_chars = 0
    for slide in presentation.slides:
        total_images += sum(1 for shape in slide.shapes if getattr(shape, "shape_type", None) == 13)
        for shape in slide.shapes:
            text_frame = getattr(shape, "text_frame", None)
            if text_frame is not None:
                total_chars += len(str(getattr(text_frame, "text", "")).strip())
    return f"上轮基线 | {slide_count}页 | {total_images}张图 | {total_chars}字"


def snapshot_round_artifacts(paths: list[str], round_no: int) -> list[str]:
    snapshots: list[str] = []
    for src in paths:
        src_path = Path(src).expanduser()
        if not src_path.exists() or not src_path.is_file():
            continue
        target = Path(with_round_version_suffix(str(src_path), round_no)).expanduser()
        try:
            if target.resolve() != src_path.resolve():
                target.parent.mkdir(parents=True, exist_ok=True)
                shutil.copy2(src_path, target)
            snapshots.append(str(target.resolve()))
        except Exception:
            try:
                snapshots.append(str(src_path.resolve()))
            except Exception:
                snapshots.append(str(src_path))
    return snapshots


def remember_office_path(metadata: dict[str, object], path: str) -> bool:
    stripped = path.strip()
    if not stripped:
        return False
    suffix = Path(stripped).suffix.lower()
    key_map = {
        ".pptx": "last_pptx_path",
        ".docx": "last_docx_path",
        ".xlsx": "last_xlsx_path",
    }
    key = key_map.get(suffix)
    if key is None:
        return False
    if metadata.get(key) == stripped:
        return False
    metadata[key] = stripped
    return True


def capture_latest_pptx(
    metadata: dict[str, object],
    *,
    roots: tuple[Path, ...],
    window_seconds: int = 180,
) -> bool:
    now = time.time()
    candidates: list[Path] = []
    for root in roots:
        try:
            if not root.exists():
                continue
        except Exception:
            continue
        try:
            for path in root.rglob("*.pptx"):
                try:
                    mtime = path.stat().st_mtime
                except Exception:
                    continue
                if now - mtime <= window_seconds:
                    candidates.append(path)
        except Exception:
            continue
    if not candidates:
        return False
    candidates.sort(key=lambda path: path.stat().st_mtime, reverse=True)
    newest = candidates[0]
    return remember_office_path(metadata, str(newest))


def is_office_edit_request(text: str) -> bool:
    query = text.lower()
    if not query:
        return False
    edit_hints = (
        "改",
        "修改",
        "调整",
        "换",
        "替换",
        "更换",
        "删",
        "删除",
        "加",
        "添加",
        "插入",
        "第一页",
        "第二页",
        "第三页",
        "单元格",
        "段落",
        "封面",
    )
    doc_hints = (
        "ppt",
        "pptx",
        "word",
        "docx",
        "excel",
        "xlsx",
        "幻灯片",
        "文档",
        "表格",
        "封面",
        "页",
    )
    return any(hint in query for hint in edit_hints) and any(hint in query for hint in doc_hints)


def is_image_generation_request(text: str) -> bool:
    query = text.lower()
    if not query:
        return False
    generation_hints = (
        "生图",
        "出图",
        "文生图",
        "图生图",
        "图像生成",
        "以图生图",
        "根据这张图生成",
        "image generation",
        "generate image",
        "text-to-image",
        "image-to-image",
        "txt2img",
        "img2img",
    )
    return any(hint in query for hint in generation_hints)


def build_image_generation_system_message() -> Message:
    return Message(
        role="system",
        content=(
            "检测到这是生图任务（文生图/图生图）。\n"
            "执行约束：\n"
            "1) 允许最多 2 次轻量探测（如检查环境变量/关键配置）；\n"
            "2) 之后必须立即写最小脚本到 /tmp 并执行一次真实请求；\n"
            "3) 禁止连续使用 ls/stat/test/echo 循环探测；\n"
            "4) 输出必须包含请求命令、HTTP 状态码、返回摘要与图片绝对路径（或明确失败原因）。"
        ),
    )


def is_complex_office_request(text: str) -> bool:
    query = text.lower()
    if not query:
        return False
    complex_hints = (
        "插图",
        "配图",
        "图片",
        "海报",
        "背景图",
        "图标",
        "视频",
        "音频",
        "音乐",
        "动效",
        "动画",
        "高端",
        "商务风",
        "版式",
        "排版",
        "重排",
        "模板",
        "封面设计",
        "视觉风格",
    )
    return any(hint in query for hint in complex_hints)


def mentions_specific_dark_bar_target(text: str) -> bool:
    query = text.lower()
    if not query:
        return False
    target_hints = ("黑色横条", "黑条", "黑色条", "深色横条", "黑色块", "黑底条")
    return any(hint in query for hint in target_hints)


def has_any_last_office_path(metadata: dict[str, object]) -> bool:
    for key in ("last_pptx_path", "last_docx_path", "last_xlsx_path"):
        value = metadata.get(key)
        if isinstance(value, str) and value.strip():
            return True
    return False


def is_followup_edit_message(text: str) -> bool:
    query = text.lower()
    if not query:
        return False
    hints = (
        "改一下",
        "修改",
        "调整",
        "换一",
        "替换",
        "更换",
        "删掉",
        "删除",
        "加一",
        "添加",
        "插入",
        "第一页",
        "第一面",
        "第二页",
        "第三页",
        "单元格",
        "段落",
        "封面图",
        "封面",
    )
    return any(hint in query for hint in hints)


def get_default_office_edit_path(
    tool_name: str,
    metadata: dict[str, object],
) -> str | None:
    key_map = {
        "ppt_edit": "last_pptx_path",
        "docx_edit": "last_docx_path",
        "xlsx_edit": "last_xlsx_path",
    }
    key = key_map.get(tool_name)
    if key is None:
        return None
    value = metadata.get(key)
    if not isinstance(value, str):
        return None
    path = value.strip()
    return path or None


def build_office_edit_hint_system_message(metadata: dict[str, object]) -> Message | None:
    hints: list[str] = []
    for label, key in (
        ("PPT", "last_pptx_path"),
        ("Word", "last_docx_path"),
        ("Excel", "last_xlsx_path"),
    ):
        value = metadata.get(key)
        if isinstance(value, str) and value.strip():
            hints.append(f"- {label}: {value.strip()}")
    if not hints:
        return None
    return Message(
        role="system",
        content=(
            "检测到用户在要求修改已有 Office 文件。\n"
            "若只是文案小改，请优先使用对应编辑工具（ppt_edit/docx_edit/xlsx_edit）；"
            "若涉及插图/音视频/风格升级/复杂排版，"
            "请先输出简短执行计划，"
            "再组合 browser、bash、file_write 与编辑工具执行。\n"
            "图片操作规则：\n"
            "- 换图/替换图片 → action=replace_image（删旧图+在原位置放新图）\n"
            "- 新增图片 → action=add_image\n"
            "- 删除图片 → action=remove_image\n"
            "严禁用 add_image 来「换图」，那样旧图还在，新图只是盖在上面。\n"
            "必须优先修改用户明确点名的对象（页码/元素/文案），"
            "不要改成泛化动作（例如只改整页背景）。\n"
            "不要把复杂请求机械降级为“只能改文字”。\n"
            "并优先使用以下最近文件路径，不要先用 bash 反复探测：\n" + "\n".join(hints)
        ),
    )


def build_office_path_block_message(metadata: dict[str, object]) -> str:
    hints: list[str] = []
    for label, key in (
        ("PPT", "last_pptx_path"),
        ("Word", "last_docx_path"),
        ("Excel", "last_xlsx_path"),
    ):
        value = metadata.get(key)
        if isinstance(value, str) and value.strip():
            hints.append(f"{label}: {value.strip()}")
    if not hints:
        return "检测到这是 Office 修改请求，请直接使用已有文件路径，不要用 bash 查找。"
    return (
        "检测到已有 Office 文件路径，禁止用 bash 探测/查找。\n"
        "请直接使用以下路径：\n" + "\n".join(f"- {hint}" for hint in hints)
    )


def build_complex_office_plan_system_message() -> Message:
    return Message(
        role="system",
        content=(
            "检测到这是复杂 Office 修改任务（可能包含插图/音视频/风格升级/版式调整）。\n"
            "执行要求：\n"
            "1) 先给用户一句话计划（将做哪些步骤）；\n"
            "2) 再按步骤调用工具，不要只调用一次 ppt_edit 就结束；\n"
            "3) 如需媒体素材，优先先获取素材文件路径，再写入文档。"
        ),
    )


def append_office_system_hints(
    system_messages: list[Message],
    session: Session | None,
    llm_message: str,
) -> None:
    if session is None:
        return
    is_office_request_now = is_office_edit_request(llm_message) or (
        is_followup_edit_message(llm_message) and has_any_last_office_path(session.metadata)
    )
    if not is_office_request_now:
        return
    office_hint = build_office_edit_hint_system_message(session.metadata)
    if office_hint is not None:
        system_messages.append(office_hint)
    if is_complex_office_request(llm_message):
        system_messages.append(build_complex_office_plan_system_message())


def force_include_office_edit_tools(
    selected: set[str],
    *,
    available: set[str],
    session: Session | None,
    llm_message: str,
) -> set[str]:
    if session is None:
        return selected
    if not (
        is_office_edit_request(llm_message)
        or (is_followup_edit_message(llm_message) and has_any_last_office_path(session.metadata))
    ):
        return selected

    key_by_tool = {
        "ppt_edit": "last_pptx_path",
        "docx_edit": "last_docx_path",
        "xlsx_edit": "last_xlsx_path",
    }
    expanded = set(selected)
    for tool_name, meta_key in key_by_tool.items():
        if tool_name not in available:
            continue
        path_value = session.metadata.get(meta_key)
        if isinstance(path_value, str) and path_value.strip():
            expanded.add(tool_name)
    return expanded


__all__ = [
    "ABS_FILE_PATH_RE",
    "NON_DELIVERY_EXTS",
    "OFFICE_PATH_RE",
    "append_office_system_hints",
    "baseline_from_pptx",
    "build_complex_office_plan_system_message",
    "build_image_generation_system_message",
    "build_office_edit_hint_system_message",
    "build_office_path_block_message",
    "capture_latest_pptx",
    "extract_artifact_baseline",
    "extract_delivery_artifact_paths",
    "extract_office_paths",
    "extract_round_delivery_section",
    "fix_version_suffix",
    "force_include_office_edit_tools",
    "get_default_office_edit_path",
    "has_any_last_office_path",
    "is_complex_office_request",
    "is_followup_edit_message",
    "is_image_generation_request",
    "is_office_edit_request",
    "is_office_path_probe_command",
    "looks_like_ppt_generation_command",
    "looks_like_ppt_generation_script",
    "mentions_specific_dark_bar_target",
    "remember_office_path",
    "snapshot_round_artifacts",
    "with_round_version_suffix",
]
