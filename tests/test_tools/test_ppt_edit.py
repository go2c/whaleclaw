"""Tests for ppt_edit tool."""

from __future__ import annotations

import base64

import pytest

from whaleclaw.tools.ppt_edit import PptEditTool

_PNG_1X1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO7Zx7kAAAAASUVORK5CYII="
)


@pytest.mark.asyncio
async def test_ppt_edit_replace_text(tmp_path) -> None:
    from pptx import Presentation

    ppt_path = tmp_path / "demo.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "旧标题"
    prs.save(str(ppt_path))

    tool = PptEditTool()
    result = await tool.execute(
        path=str(ppt_path),
        slide_index=1,
        old_text="旧标题",
        new_text="新标题",
    )
    assert result.success is True

    prs2 = Presentation(str(ppt_path))
    assert prs2.slides[0].shapes.title.text == "新标题"


@pytest.mark.asyncio
async def test_ppt_edit_set_notes_and_background(tmp_path) -> None:
    from pptx import Presentation

    ppt_path = tmp_path / "style.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(str(ppt_path))

    tool = PptEditTool()
    r1 = await tool.execute(
        path=str(ppt_path),
        slide_index=1,
        action="set_notes",
        new_text="这是备注",
    )
    r2 = await tool.execute(
        path=str(ppt_path),
        slide_index=1,
        action="set_background",
        background_color="#0F2747",
    )
    assert r1.success is True
    assert r2.success is True

    prs2 = Presentation(str(ppt_path))
    notes_text = prs2.slides[0].notes_slide.notes_text_frame.text
    assert "这是备注" in notes_text


@pytest.mark.asyncio
async def test_ppt_edit_add_image(tmp_path) -> None:
    from pptx import Presentation

    ppt_path = tmp_path / "image.pptx"
    img_path = tmp_path / "dot.png"
    img_path.write_bytes(_PNG_1X1)

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(ppt_path))

    tool = PptEditTool()
    result = await tool.execute(
        path=str(ppt_path),
        slide_index=1,
        action="add_image",
        image_path=str(img_path),
        image_width=2.0,
    )
    assert result.success is True

    prs2 = Presentation(str(ppt_path))
    assert len(prs2.slides[0].shapes) >= 1


@pytest.mark.asyncio
async def test_ppt_edit_apply_business_style_restyles_dark_bar(tmp_path) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    ppt_path = tmp_path / "bar_style.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(
        1,  # rectangle
        Inches(0),
        Inches(0.8),
        Inches(12.8),
        Inches(2.0),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0, 0, 0)
    prs.save(str(ppt_path))

    tool = PptEditTool()
    result = await tool.execute(
        path=str(ppt_path),
        slide_index=1,
        action="apply_business_style",
    )
    assert result.success is True
    assert "重设深色条 1 处" in result.output

    prs2 = Presentation(str(ppt_path))
    fill_rgb = prs2.slides[0].shapes[0].fill.fore_color.rgb
    assert fill_rgb[0] > 0 and fill_rgb[1] > 0 and fill_rgb[2] > 0
