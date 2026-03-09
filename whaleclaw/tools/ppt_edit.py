"""PPT edit tool — edit text/style/image/notes in an existing .pptx file."""

from __future__ import annotations

from pathlib import Path
from typing import Any, cast

from whaleclaw.tools.base import Tool, ToolDefinition, ToolParameter, ToolResult
from whaleclaw.tools.deps import ensure_tool_dep


def _add_cropped_picture(
    slide: Any,
    img_path: Path,
    left: int,
    top: int,
    target_w: int,
    target_h: int,
) -> Any:
    """Add a picture to *slide* with smart face-aware cropping.

    Returns the created picture shape, or None if fallback (no PIL) was used.
    """
    try:
        from PIL import Image as PILImage

        from whaleclaw.utils.image_crop import detect_face_info, smart_crop_box

        with PILImage.open(str(img_path)) as im:
            iw, ih = im.size

        fi = detect_face_info(str(img_path))
        crop_box = smart_crop_box(iw, ih, iw, int(iw * target_h / target_w), face_info=fi)
        cx0, cy0, cx1, cy1 = crop_box

        crop_left_frac = cx0 / iw
        crop_right_frac = 1.0 - cx1 / iw
        crop_top_frac = cy0 / ih
        crop_bottom_frac = 1.0 - cy1 / ih

        img_ratio = iw / ih
        box_ratio = target_w / target_h
        if img_ratio > box_ratio:
            scale_h = target_h
            scale_w = int(target_h * img_ratio)
        else:
            scale_w = target_w
            scale_h = int(target_w / img_ratio)

        pic = slide.shapes.add_picture(
            str(img_path),
            int(left - (scale_w - target_w) / 2),
            int(top - (scale_h - target_h) / 2),
            int(scale_w),
            int(scale_h),
        )
        pic.crop_left = crop_left_frac
        pic.crop_right = crop_right_frac
        pic.crop_top = crop_top_frac
        pic.crop_bottom = crop_bottom_frac
        pic.left = int(left)
        pic.top = int(top)
        pic.width = int(target_w)
        pic.height = int(target_h)
        return pic
    except ImportError:
        slide.shapes.add_picture(str(img_path), left, top, width=target_w)
        return None


class PptEditTool(Tool):
    """Edit an existing PPT by actions on a specific slide."""

    @property
    def definition(self) -> ToolDefinition:
        return ToolDefinition(
            name="ppt_edit",
            description=(
                "修改现有 PPT（.pptx）：支持文本替换、标题更新、商务风格、背景色、备注、"
                "插图（add_image 新增图片；replace_image 替换已有图片，保持原位置尺寸；"
                "remove_image 删除图片）。换图/替换封面图请用 replace_image。"
            ),
            parameters=[
                ToolParameter(name="path", type="string", description="PPT 文件绝对路径"),
                ToolParameter(name="slide_index", type="integer", description="页码（从 1 开始）"),
                ToolParameter(
                    name="action",
                    type="string",
                    description=(
                        "操作类型：replace_text|set_title|set_notes|set_background"
                        "|add_image|replace_image|remove_image|apply_business_style"
                    ),
                    required=False,
                    enum=[
                        "replace_text",
                        "set_title",
                        "set_notes",
                        "set_background",
                        "add_image",
                        "replace_image",
                        "remove_image",
                        "apply_business_style",
                    ],
                ),
                ToolParameter(
                    name="image_index",
                    type="integer",
                    description=(
                        "目标图片序号（从 1 开始，replace_image/remove_image 时使用，"
                        "默认 1 即第一张图）"
                    ),
                    required=False,
                ),
                ToolParameter(
                    name="old_text",
                    type="string",
                    description="旧文案（replace_text 时必填）",
                    required=False,
                ),
                ToolParameter(
                    name="new_text",
                    type="string",
                    description="新文案（replace_text/set_title/set_notes 时使用）",
                    required=False,
                ),
                ToolParameter(
                    name="image_path",
                    type="string",
                    description="插图绝对路径（add_image 时必填）",
                    required=False,
                ),
                ToolParameter(
                    name="image_left",
                    type="number",
                    description="插图左边距（英寸，默认 1.0）",
                    required=False,
                ),
                ToolParameter(
                    name="image_top",
                    type="number",
                    description="插图上边距（英寸，默认 1.8）",
                    required=False,
                ),
                ToolParameter(
                    name="image_width",
                    type="number",
                    description="插图宽度（英寸，默认 6.5）",
                    required=False,
                ),
                ToolParameter(
                    name="image_height",
                    type="number",
                    description=(
                        "插图高度（英寸）。同时指定 width+height 时自动裁剪保持比例，"
                        "不会变形拉伸。"
                    ),
                    required=False,
                ),
                ToolParameter(
                    name="background_color",
                    type="string",
                    description="背景色（十六进制，如 #0F2747）",
                    required=False,
                ),
            ],
        )

    async def execute(self, **kwargs: Any) -> ToolResult:
        if not ensure_tool_dep("pptx"):
            return ToolResult(success=False, output="", error="缺少依赖 python-pptx")

        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        raw_path = str(kwargs.get("path", "")).strip()
        action = str(kwargs.get("action", "replace_text")).strip().lower() or "replace_text"
        old_text = str(kwargs.get("old_text", "")).strip()
        new_text = str(kwargs.get("new_text", ""))
        image_path = str(kwargs.get("image_path", "")).strip()
        background_color = str(kwargs.get("background_color", "")).strip()
        try:
            slide_index = int(kwargs.get("slide_index", 0))
        except (TypeError, ValueError):
            slide_index = 0

        if not raw_path:
            return ToolResult(success=False, output="", error="path 不能为空")
        if slide_index <= 0:
            return ToolResult(success=False, output="", error="slide_index 必须 >= 1")

        path = Path(raw_path).expanduser().resolve()
        if not path.is_file():
            return ToolResult(success=False, output="", error=f"文件不存在: {path}")
        if path.suffix.lower() != ".pptx":
            return ToolResult(success=False, output="", error="仅支持 .pptx 文件")

        try:
            prs = Presentation(str(path))
        except Exception as exc:
            return ToolResult(success=False, output="", error=f"PPT 打开失败: {exc}")

        if slide_index > len(prs.slides):
            return ToolResult(
                success=False,
                output="",
                error=f"页码越界: slide_index={slide_index}, 总页数={len(prs.slides)}",
            )

        slide = prs.slides[slide_index - 1]

        if action == "replace_text":
            if not old_text:
                return ToolResult(success=False, output="", error="old_text 不能为空")
            replaced = 0
            for shape in slide.shapes:
                text = getattr(shape, "text", None)
                if not isinstance(text, str):
                    continue
                count = text.count(old_text)
                if count <= 0:
                    continue
                cast(Any, shape).text = text.replace(old_text, new_text)
                replaced += count
            if replaced == 0:
                return ToolResult(
                    success=False,
                    output="",
                    error=f"第 {slide_index} 页未找到匹配文本",
                )
            output = f"已修改 {path} 第 {slide_index} 页，替换 {replaced} 处文本"
        elif action == "set_title":
            title_shape = slide.shapes.title
            if title_shape is None:
                title_shape = slide.shapes.add_textbox(
                    Inches(0.8),
                    Inches(0.35),
                    Inches(11.0),
                    Inches(1.0),
                )
            tf = title_shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = new_text
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(34)
                run.font.name = "Microsoft YaHei"
                run.font.color.rgb = RGBColor(15, 39, 71)
            output = f"已更新 {path} 第 {slide_index} 页标题"
        elif action == "set_notes":
            notes = slide.notes_slide.notes_text_frame
            if notes is None:
                return ToolResult(success=False, output="", error="当前页备注区域不可用")
            notes.clear()
            notes.text = new_text
            output = f"已更新 {path} 第 {slide_index} 页备注"
        elif action == "set_background":
            hex_color = (background_color or "#F4F7FB").strip().lstrip("#")
            if len(hex_color) != 6:
                return ToolResult(
                    success=False,
                    output="",
                    error="background_color 必须是 6 位十六进制",
                )
            try:
                r = int(hex_color[0:2], 16)
                g = int(hex_color[2:4], 16)
                b = int(hex_color[4:6], 16)
            except ValueError:
                return ToolResult(
                    success=False,
                    output="",
                    error="background_color 不是合法十六进制颜色",
                )
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(r, g, b)
            output = f"已设置 {path} 第 {slide_index} 页背景色 #{hex_color.upper()}"
        elif action == "add_image":
            if not image_path:
                return ToolResult(success=False, output="", error="add_image 需要 image_path")
            img = Path(image_path).expanduser().resolve()
            if not img.is_file():
                return ToolResult(success=False, output="", error=f"图片不存在: {img}")
            left = Inches(float(kwargs.get("image_left", 1.0)))
            top = Inches(float(kwargs.get("image_top", 1.8)))
            target_w = Inches(float(kwargs.get("image_width", 6.5)))
            raw_h = kwargs.get("image_height")
            if raw_h is not None:
                target_h = Inches(float(raw_h))
                _add_cropped_picture(slide, img, int(left), int(top), int(target_w), int(target_h))
            else:
                slide.shapes.add_picture(str(img), left, top, width=target_w)
            output = f"已在 {path} 第 {slide_index} 页插入图片 {img.name}"
        elif action in ("replace_image", "remove_image"):
            image_shapes = [
                s for s in slide.shapes
                if s.shape_type == MSO_SHAPE_TYPE.PICTURE
            ]
            if not image_shapes:
                return ToolResult(
                    success=False, output="",
                    error=f"第 {slide_index} 页没有图片可操作",
                )
            try:
                img_idx = int(kwargs.get("image_index", 1))
            except (TypeError, ValueError):
                img_idx = 1
            if img_idx < 1 or img_idx > len(image_shapes):
                return ToolResult(
                    success=False, output="",
                    error=f"image_index={img_idx} 越界，该页共 {len(image_shapes)} 张图片",
                )
            old_shape = image_shapes[img_idx - 1]
            old_left, old_top = old_shape.left, old_shape.top
            old_width, old_height = old_shape.width, old_shape.height

            sp_elem = old_shape._element  # pyright: ignore[reportPrivateUsage]
            sp_elem_any = cast(Any, sp_elem)
            parent = sp_elem_any.getparent()
            if parent is None:
                return ToolResult(success=False, output="", error="无法定位图片父节点")
            parent.remove(sp_elem_any)

            if action == "remove_image":
                output = f"已删除 {path} 第 {slide_index} 页第 {img_idx} 张图片"
            else:
                if not image_path:
                    return ToolResult(
                        success=False, output="",
                        error="replace_image 需要 image_path",
                    )
                img = Path(image_path).expanduser().resolve()
                if not img.is_file():
                    return ToolResult(success=False, output="", error=f"图片不存在: {img}")
                _add_cropped_picture(
                    slide, img, old_left, old_top, old_width, old_height,
                )
                output = f"已替换 {path} 第 {slide_index} 页第 {img_idx} 张图片为 {img.name}"
        elif action == "apply_business_style":
            def _iter_shapes_recursive(container: Any) -> list[Any]:
                items: list[Any] = []
                for shp in container:
                    items.append(shp)
                    if getattr(shp, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                        inner = getattr(shp, "shapes", None)
                        if inner is not None:
                            items.extend(_iter_shapes_recursive(inner))
                return items

            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(245, 248, 252)
            dark_bar_restyled = 0
            for shape in _iter_shapes_recursive(slide.shapes):
                if not hasattr(shape, "fill"):
                    continue
                shape_fill = shape.fill
                if getattr(shape_fill, "type", None) != 1:  # MSO_FILL.SOLID
                    continue
                color = getattr(shape_fill, "fore_color", None)
                rgb = getattr(color, "rgb", None)
                if rgb is None:
                    continue
                shape_width = getattr(shape, "width", None)
                shape_height = getattr(shape, "height", None)
                if not isinstance(shape_width, int) or not isinstance(shape_height, int):
                    continue
                # Treat large dark rectangles/bars as candidate overlays.
                is_dark = rgb[0] <= 50 and rgb[1] <= 50 and rgb[2] <= 50
                is_large_bar = (
                    shape_width >= int(cast(int, prs.slide_width) * 0.35)
                    and shape_height <= int(cast(int, prs.slide_height) * 0.75)
                )
                if not (is_dark and is_large_bar):
                    continue
                shape_fill.solid()
                shape_fill.fore_color.rgb = RGBColor(214, 228, 245)
                if hasattr(shape_fill, "transparency"):
                    shape_fill.transparency = 0.15
                if hasattr(shape, "line") and getattr(shape.line, "fill", None) is not None:
                    shape.line.fill.background()
                dark_bar_restyled += 1
            if slide.shapes.title is not None:
                ttf = slide.shapes.title.text_frame
                if ttf.paragraphs:
                    p = ttf.paragraphs[0]
                    for run in p.runs:
                        run.font.bold = True
                        run.font.size = Pt(34)
                        run.font.name = "Microsoft YaHei"
                        run.font.color.rgb = RGBColor(15, 39, 71)
            output = (
                f"已应用 {path} 第 {slide_index} 页商务风格，"
                f"重设深色条 {dark_bar_restyled} 处"
            )
        else:
            return ToolResult(success=False, output="", error=f"不支持的 action: {action}")

        sw = cast(int, prs.slide_width)
        sh = cast(int, prs.slide_height)
        for s in slide.shapes:
            width_type = type(s.width)
            height_type = type(s.height)
            overflow_r = s.left + s.width - sw
            overflow_b = s.top + s.height - sh
            if overflow_r > 0 or overflow_b > 0:
                is_pic = getattr(s, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE
                if is_pic and s.width > 0 and s.height > 0:
                    ratio = s.width / s.height
                    if overflow_r > 0:
                        s.width = width_type(sw - s.left)
                        s.height = height_type(int(s.width / ratio))
                    if s.top + s.height > sh:
                        s.height = height_type(sh - s.top)
                        s.width = width_type(int(s.height * ratio))
                else:
                    if overflow_r > 0:
                        s.width = width_type(sw - s.left)
                    if overflow_b > 0:
                        s.height = height_type(sh - s.top)

        try:
            prs.save(str(path))
        except Exception as exc:
            return ToolResult(success=False, output="", error=f"PPT 保存失败: {exc}")

        return ToolResult(success=True, output=output)
